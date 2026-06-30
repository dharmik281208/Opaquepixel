import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    WHITE_BG = RGBColor(255, 255, 255)
    TEXT_MAIN = RGBColor(15, 23, 42)       # #0f172a
    TEXT_MUTED = RGBColor(71, 85, 105)     # #475569
    BLUE_COLOR = RGBColor(37, 99, 235)     # #2563eb
    GREEN_COLOR = RGBColor(16, 185, 129)   # #10b981
    CARD_BG_LIGHT = RGBColor(248, 250, 252)
    CARD_BG_DARK = RGBColor(15, 23, 42)    # Dark background for Conclusion
    IBM_BLUE = RGBColor(0, 102, 204)

    def add_branding(slide, dark_theme=False):
        # Top Left: Bharat Cares Logo
        txBox_left = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(3.5), Inches(0.6))
        tf_l = txBox_left.text_frame
        p_l = tf_l.paragraphs[0]
        p_l.text = "Bharat Cares"
        p_l.font.size = Pt(20)
        p_l.font.bold = True
        p_l.font.color.rgb = RGBColor(0, 168, 225)
        p_sub = tf_l.add_paragraph()
        p_sub.text = "by SMEC Trust"
        p_sub.font.size = Pt(8.5)
        p_sub.font.color.rgb = RGBColor(203, 213, 225) if dark_theme else TEXT_MUTED

        # Top Right: IBM Collaboration
        txBox_right = slide.shapes.add_textbox(Inches(10.2), Inches(0.2), Inches(2.5), Inches(0.6))
        tf_r = txBox_right.text_frame
        p_r = tf_r.paragraphs[0]
        p_r.text = "In collaboration with"
        p_r.font.size = Pt(8.5)
        p_r.font.color.rgb = RGBColor(203, 213, 225) if dark_theme else TEXT_MUTED
        p_r.alignment = PP_ALIGN.RIGHT
        
        p_ibm = tf_r.add_paragraph()
        p_ibm.text = "IBM"
        p_ibm.font.size = Pt(24)
        p_ibm.font.bold = True
        p_ibm.font.color.rgb = RGBColor(96, 165, 250) if dark_theme else IBM_BLUE
        p_ibm.alignment = PP_ALIGN.RIGHT

    # ==========================================
    # SLIDE 1: What Did You Learn in This Program?
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = WHITE_BG
    bg1.line.fill.background()
    add_branding(slide1)

    txBox1 = slide1.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(1.2))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "What Did You Learn in This Program?"
    p1.font.size = Pt(30)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    p1.space_after = Pt(4)

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Through this program, I gained a foundational introduction to steganography & cybersecurity, built core technical skills, and understood the critical importance of security in enterprise operations — learning to analyze and mitigate real-world vulnerabilities through hands-on engineering."
    p1_sub.font.size = Pt(12.5)
    p1_sub.font.italic = True
    p1_sub.font.color.rgb = TEXT_MUTED

    learnings = [
        ("1", BLUE_COLOR, "Multi-Domain Steganography & Frequency Transform Engineering",
         "Implementing Discrete Cosine Transform (DCT) & Discrete Sine Transform (DST) frequency-domain embedding, spatial Least Significant Bit (LSB) manipulation, and Pixel Value Differencing (PVD) edge variance algorithms across image, audio (WAV), video (MP4), and document (PDF) carriers."),
        ("2", GREEN_COLOR, "Statistical Stegananalysis & Universal Anomaly Inspection",
         "Developing automated Chi-Square (χ²) sample pairs analysis, Shannon entropy scoring (0.0–8.0), and high-frequency coefficient residual scanning to detect third-party hidden payloads and structural container tampering."),
        ("3", BLUE_COLOR, "High-Performance Decoupled Microservice Architecture",
         "Building an asynchronous FastAPI Python 3.14 REST API backend integrated with NumPy and OpenCV (cv2) matrix processing pipelines, engineered with automated CORS policies and stateless REST session tokens."),
        ("4", GREEN_COLOR, "Modern Responsive UI Engineering & Liquid Glass Design System",
         "Crafting an immersive macOS Tahoe Liquid Glass user interface with Vite, React 18, and Tailwind CSS, featuring an open-access frictionless workflow, dynamic dark/light theme switching, and live interactive scanning overlays."),
        ("5", BLUE_COLOR, "Production CI/CD Automation & Distributed Cloud Deployment",
         "Engineering automated GitHub Actions workflows for production bundle compilation, managing single-page application (SPA) routing fallbacks (404.html), and hosting a distributed ecosystem across GitHub Pages (opaquepixel.app) and Render Cloud Web Services.")
    ]

    top_pos = 2.4
    for num, color, title, desc in learnings:
        circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(top_pos), Inches(0.38), Inches(0.38))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf_c = circle.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = num
        p_c.font.size = Pt(13)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE_BG
        p_c.alignment = PP_ALIGN.CENTER

        txBox_item = slide1.shapes.add_textbox(Inches(1.1), Inches(top_pos - 0.08), Inches(11.5), Inches(0.85))
        tf_i = txBox_item.text_frame
        tf_i.word_wrap = True
        p_t = tf_i.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MAIN
        p_t.space_after = Pt(2)
        p_d = tf_i.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        top_pos += 0.95

    # ==========================================
    # SLIDE 2: Project Purpose (2x2 Grid)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = WHITE_BG
    bg2.line.fill.background()
    add_branding(slide2)

    txBox2 = slide2.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.9))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Project Purpose"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MAIN
    p2.space_after = Pt(2)
    p2_sub = tf2.add_paragraph()
    p2_sub.text = "Why Did You Make This Project?"
    p2_sub.font.size = Pt(16)
    p2_sub.font.italic = True
    p2_sub.font.color.rgb = TEXT_MUTED

    purposes = [
        ("1", BLUE_COLOR, "Vulnerability of Unencrypted Communication",
         "Traditional file transfers expose plain data to network surveillance and unauthorized interception — OpaquePixel hides secret encrypted payloads inside ordinary multimedia files to ensure undetectable transmission.", 0.6, 2.0),
        ("2", GREEN_COLOR, "Need for Compression-Resilient Steganography",
         "Standard spatial LSB methods fail under modern messaging compression — OpaquePixel implements Discrete Cosine Transform (DCT) and F5 matrix algorithms to maintain data integrity across compressed channels.", 6.8, 2.0),
        ("3", BLUE_COLOR, "Detection of Covert Steganographic Threats",
         "Security analysts face difficulty identifying hidden data embedded by malicious tools — OpaquePixel provides automated Chi-Square (χ²) steganalysis and entropy probes to detect covert payloads and container anomalies.", 0.6, 4.6),
        ("4", GREEN_COLOR, "Unified Forensic Dashboard for Security Teams",
         "Eliminates complex command-line overhead by consolidating multi-algorithm encoding, real-time forensic scanning, and automated PDF report generation into a single glassmorphism web platform (opaquepixel.app).", 6.8, 4.6)
    ]

    for num, color, title, desc, left, top in purposes:
        card_shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.9), Inches(2.2))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = CARD_BG_LIGHT
        card_shape.line.fill.background()

        circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.25), Inches(top + 0.25), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf_c = circle.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = num
        p_c.font.size = Pt(15)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE_BG
        p_c.alignment = PP_ALIGN.CENTER

        txBox_card = slide2.shapes.add_textbox(Inches(left + 0.8), Inches(top + 0.18), Inches(4.9), Inches(1.8))
        tf_card = txBox_card.text_frame
        tf_card.word_wrap = True
        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MAIN
        p_t.space_after = Pt(3)
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 3: Project Objectives (5 List Cards)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = WHITE_BG
    bg3.line.fill.background()
    add_branding(slide3)

    txBox3 = slide3.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.9))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Project Objectives"
    p3.font.size = Pt(30)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_MAIN
    p3.space_after = Pt(2)
    p3_sub = tf3.add_paragraph()
    p3_sub.text = "What Did You Make in This Project?"
    p3_sub.font.size = Pt(16)
    p3_sub.font.italic = True
    p3_sub.font.color.rgb = TEXT_MUTED

    objectives = [
        (BLUE_COLOR, "Multi-Carrier Steganography Engine", "Asynchronous Python backend (routers/hide.py, routers/reveal.py) embedding encrypted payloads inside PNG, JPG, MP4, WAV, and PDF carriers."),
        (GREEN_COLOR, "F5 Matrix & Frequency Domain Encoding Core", "Advanced spatial and spectral algorithms (F5, DST, LSB, PVD) engineered for compression resilience and high payload security."),
        (BLUE_COLOR, "Automated Stegananalysis & Forensic Audit Probe", "Real-time scanning service (routers/scan.py) evaluating Chi-Square (χ²) bit-plane anomalies, Shannon entropy (0.0–8.0), and coefficient perturbations."),
        (GREEN_COLOR, "PDF Forensic Report Generation Utility", "Client-side reporting tool (pdfReport.js) synthesizing multi-algorithm audit scores into downloadable, printable forensic compliance reports."),
        (BLUE_COLOR, "macOS Tahoe Liquid Glass Workspace UI", "Modern responsive React 18 interface (opaquepixel.app) featuring open-access workflows, dynamic dark/light theme switching, and live interactive scanning overlays.")
    ]

    top_pos3 = 2.0
    for color, title, desc in objectives:
        card_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(top_pos3), Inches(12.1), Inches(0.95))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = CARD_BG_LIGHT
        card_shape.line.fill.background()

        dot = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(top_pos3 + 0.36), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        txBox_card = slide3.shapes.add_textbox(Inches(1.3), Inches(top_pos3 + 0.12), Inches(11.2), Inches(0.75))
        tf_card = txBox_card.text_frame
        tf_card.word_wrap = True
        p_item = tf_card.paragraphs[0]
        
        run_t = p_item.add_run()
        run_t.text = title + " — "
        run_t.font.size = Pt(13.5)
        run_t.font.bold = True
        run_t.font.color.rgb = TEXT_MAIN

        run_d = p_item.add_run()
        run_d.text = desc
        run_d.font.size = Pt(12.5)
        run_d.font.color.rgb = TEXT_MUTED

        top_pos3 += 1.05

    # ==========================================
    # SLIDE 4: Structure of Your Project (2x3 Grid Dark Cards)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = WHITE_BG
    bg4.line.fill.background()
    add_branding(slide4)

    txBox4 = slide4.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.9))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Structure of Your Project"
    p4.font.size = Pt(30)
    p4.font.bold = True
    p4.font.color.rgb = TEXT_MAIN
    p4.space_after = Pt(2)
    p4_sub = tf4.add_paragraph()
    p4_sub.text = "What Are the Steps You Followed?"
    p4_sub.font.size = Pt(16)
    p4_sub.font.italic = True
    p4_sub.font.color.rgb = TEXT_MUTED

    steps_data = [
        ("01", GREEN_COLOR, "Media Carrier Selection & Payload Prep",
         "Client-side validation of target carrier format (PNG, JPG, MP4, WAV, PDF) and AES-256 payload encryption parameter initialization.", 0.6, 2.0),
        ("02", GREEN_COLOR, "REST API Microservice Dispatch",
         "Asynchronous POST submission (/api/hide, /api/reveal) transmitting multipart form data to the FastAPI Python backend.", 4.7, 2.0),
        ("03", GREEN_COLOR, "Multi-Algorithm Steganography Core",
         "Executing spatial (LSB, PVD) or frequency-domain (F5 Matrix, DST) embedding engines to modify coefficient matrices without visual distortion.", 8.8, 2.0),
        ("04", GREEN_COLOR, "Automated Stegananalysis & Forensic Probe",
         "Real-time carrier inspection (/api/scan) performing Chi-Square (χ²) bit-plane auditing, Shannon entropy scoring, and container anomaly checks.", 0.6, 4.6),
        ("05", GREEN_COLOR, "Client-Side PDF Audit Report Generation",
         "Synthesizing forensic evaluation scores and container metadata into downloadable, printable compliance reports (pdfReport.js).", 4.7, 4.6),
        ("06", GREEN_COLOR, "Distributed CI/CD Cloud Deployment",
         "Automated GitHub Actions production compilation hosting frontend static SPA routes on GitHub Pages (opaquepixel.app) and backend APIs on Render.", 8.8, 4.6)
    ]

    for num, num_color, title, desc, left, top in steps_data:
        card_shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.9), Inches(2.2))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = CARD_BG_DARK
        card_shape.line.fill.background()

        txBox_card = slide4.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(3.5), Inches(1.9))
        tf_card = txBox_card.text_frame
        tf_card.word_wrap = True
        
        p_n = tf_card.paragraphs[0]
        p_n.text = num
        p_n.font.size = Pt(22)
        p_n.font.bold = True
        p_n.font.color.rgb = num_color
        p_n.space_after = Pt(4)

        p_t = tf_card.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE_BG
        p_t.space_after = Pt(3)

        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = RGBColor(203, 213, 225)

    # ==========================================
    # SLIDE 5: Advantages and Limitations (Side-by-Side Banners)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = WHITE_BG
    bg5.line.fill.background()
    add_branding(slide5)

    txBox5 = slide5.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.8))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Advantages and Limitations of Your Project"
    p5.font.size = Pt(30)
    p5.font.bold = True
    p5.font.color.rgb = TEXT_MAIN

    banner_left = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(5.9), Inches(0.6))
    banner_left.fill.solid()
    banner_left.fill.fore_color.rgb = GREEN_COLOR
    banner_left.line.fill.background()
    tf_bl = banner_left.text_frame
    p_bl = tf_bl.paragraphs[0]
    p_bl.text = "ADVANTAGES"
    p_bl.font.size = Pt(16)
    p_bl.font.bold = True
    p_bl.font.color.rgb = WHITE_BG
    p_bl.alignment = PP_ALIGN.CENTER

    card_left = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.5), Inches(5.9), Inches(4.5))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG_LIGHT
    card_left.line.fill.background()

    txBox_adv = slide5.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(5.5), Inches(4.1))
    tf_adv = txBox_adv.text_frame
    tf_adv.word_wrap = True

    adv_items = [
        "High Payload Security & Compression Resilience (F5 Matrix & DST Encoding)",
        "Universal Stegananalysis Probe (Chi-Square χ² & Entropy Auditing)",
        "Multi-Carrier Format Flexibility (PNG, JPG, WAV, MP4, PDF)",
        "Automated Compliance PDF Audit Report Generation",
        "Open-Access Frictionless UX with Dual Theme Engine (Dark & White)",
        "Automated Continuous Cloud Deployment (GitHub Pages & Render Services)"
    ]

    for idx, item in enumerate(adv_items):
        p_adv = tf_adv.paragraphs[0] if idx == 0 else tf_adv.add_paragraph()
        p_adv.text = "✓   " + item
        p_adv.font.size = Pt(12)
        p_adv.font.color.rgb = TEXT_MAIN
        p_adv.space_after = Pt(14)

    banner_right = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.9), Inches(0.6))
    banner_right.fill.solid()
    banner_right.fill.fore_color.rgb = BLUE_COLOR
    banner_right.line.fill.background()
    tf_br = banner_right.text_frame
    p_br = tf_br.paragraphs[0]
    p_br.text = "LIMITATIONS"
    p_br.font.size = Pt(16)
    p_br.font.bold = True
    p_br.font.color.rgb = WHITE_BG
    p_br.alignment = PP_ALIGN.CENTER

    card_right = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.5), Inches(5.9), Inches(4.5))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG_LIGHT
    card_right.line.fill.background()

    txBox_lim = slide5.shapes.add_textbox(Inches(7.0), Inches(2.7), Inches(5.5), Inches(4.1))
    tf_lim = txBox_lim.text_frame
    tf_lim.word_wrap = True

    lim_items = [
        "Extremely destructive image re-compression (e.g. heavy WhatsApp compression) can strip LSB payload bits",
        "Large video and high-resolution carrier uploads require higher network bandwidth",
        "Free tier cloud web services (Render) experience cold-start latency after inactivity"
    ]

    for idx, item in enumerate(lim_items):
        p_lim = tf_lim.paragraphs[0] if idx == 0 else tf_lim.add_paragraph()
        p_lim.text = "✕   " + item
        p_lim.font.size = Pt(12)
        p_lim.font.color.rgb = TEXT_MAIN
        p_lim.space_after = Pt(18)

    # ==========================================
    # SLIDE 6: Challenges You Faced (2x2 Grid with Green Resolutions)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = WHITE_BG
    bg6.line.fill.background()
    add_branding(slide6)

    txBox6 = slide6.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.9))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Challenges You Faced"
    p6.font.size = Pt(30)
    p6.font.bold = True
    p6.font.color.rgb = TEXT_MAIN
    p6.space_after = Pt(2)
    p6_sub = tf6.add_paragraph()
    p6_sub.text = "What Difficulties Did You Face?"
    p6_sub.font.size = Pt(16)
    p6_sub.font.italic = True
    p6_sub.font.color.rgb = TEXT_MUTED

    challenges_data = [
        ("Matrix Coefficient Perturbation in Compressed Images",
         "Preventing payload corruption when carrier images undergo lossy JPEG quantization or DCT compression.",
         "Resolved by engineering F5 Matrix permutation and DST frequency-domain Quantization Index Modulation (QIM).", 0.6, 2.0),
        
        ("Detecting Third-Party Steganography without Keys",
         "Auditing carrier images for hidden data inserted by external tools without having access to the encryption password.",
         "Resolved using universal Chi-Square (χ²) sample pair distribution analysis and Shannon entropy scoring.", 6.8, 2.0),
        
        ("SPA Routing & Fallbacks on GitHub Pages",
         "Eliminating HTTP 404 File Not Found errors when users reload deep routes or direct links on static cloud hosting.",
         "Resolved by automating 404.html fallback page generation in GitHub Actions CI/CD workflows.", 0.6, 4.6),
        
        ("Maintaining High-Contrast Accessibility Across Themes",
         "Guaranteeing complete legibility for dark bioluminescent and monochrome white SaaS themes without structural shifts.",
         "Resolved by standardizing unified box-model geometry and dynamic contrast CSS variables.", 6.8, 4.6)
    ]

    for title, desc, res, left, top in challenges_data:
        card_shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.9), Inches(2.2))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = CARD_BG_LIGHT
        card_shape.line.fill.background()

        txBox_card = slide6.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.18), Inches(5.3), Inches(1.8))
        tf_card = txBox_card.text_frame
        tf_card.word_wrap = True

        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MAIN
        p_t.space_after = Pt(3)

        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(6)

        p_r = tf_card.add_paragraph()
        p_r.text = "✓ " + res
        p_r.font.size = Pt(11)
        p_r.font.bold = True
        p_r.font.color.rgb = GREEN_COLOR

    # ==========================================
    # SLIDE 7: Conclusion & Future Scope (Dark Theme matching reference)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = CARD_BG_DARK
    bg7.line.fill.background()
    add_branding(slide7, dark_theme=True)

    txBox7 = slide7.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.8))
    tf7 = txBox7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Conclusion"
    p7.font.size = Pt(36)
    p7.font.bold = True
    p7.font.color.rgb = WHITE_BG

    conclusions = [
        ("1", "Robust Multi-Domain Steganography Engine", "OpaquePixel successfully bridges spatial (LSB, PVD) and frequency-domain (F5 Matrix, DST) algorithms to provide undetectable, compression-resilient data embedding across images, audio, video, and documents."),
        ("2", "Automated Universal Stegananalysis Probe", "Combines real-time Chi-Square (χ²) bit-plane auditing, Shannon entropy scoring, and structural container inspection to detect covert threats and generate downloadable compliance PDF audit reports."),
        ("3", "Enterprise Cloud-Ready Ecosystem", "Delivered as an open-access, decoupled architecture with a modern macOS Tahoe Liquid Glass UI on GitHub Pages (opaquepixel.app) and a high-performance Python 3.14 FastAPI backend on Render.")
    ]

    top_pos7 = 1.9
    for num, title, desc in conclusions:
        circle = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(top_pos7), Inches(0.38), Inches(0.38))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GREEN_COLOR
        circle.line.fill.background()
        tf_c = circle.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = num
        p_c.font.size = Pt(14)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE_BG
        p_c.alignment = PP_ALIGN.CENTER

        txBox_item = slide7.shapes.add_textbox(Inches(1.1), Inches(top_pos7 - 0.08), Inches(11.5), Inches(0.85))
        tf_i = txBox_item.text_frame
        tf_i.word_wrap = True
        p_t = tf_i.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE_BG
        p_t.space_after = Pt(2)
        p_d = tf_i.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = RGBColor(203, 213, 225)
        top_pos7 += 1.05

    # Future Scope Bottom Banner Card
    fs_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.3))
    fs_card.fill.solid()
    fs_card.fill.fore_color.rgb = RGBColor(30, 41, 59) # translucent slate dark box
    fs_card.line.color.rgb = RGBColor(51, 65, 85)

    txBox_fs = slide7.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.9))
    tf_fs = txBox_fs.text_frame
    tf_fs.word_wrap = True
    p_fs = tf_fs.paragraphs[0]
    
    run_fs_label = p_fs.add_run()
    run_fs_label.text = "FUTURE SCOPE — "
    run_fs_label.font.size = Pt(13.5)
    run_fs_label.font.bold = True
    run_fs_label.font.color.rgb = GREEN_COLOR

    run_fs_text = p_fs.add_run()
    run_fs_text.text = "Deep neural network (CNN) steganalysis classification models, real-time video stream steganography, and hardware security module (HSM) cryptographic key storage integration."
    run_fs_text.font.size = Pt(13)
    run_fs_text.font.color.rgb = RGBColor(203, 213, 225)

    output_path = r"c:\Users\LENOVO\PycharmProjects\OpaquePixel\OpaquePixel_Master_Final_Presentation.pptx"
    prs.save(output_path)
    print(f"Master Final 7-Slide Presentation created successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
